"""MVP de agente RAG corporativo con embeddings locales y ChromaDB.

La aplicación no necesita claves de API: Chroma descarga una vez un modelo de
embeddings liviano y guarda el índice vectorial en data/chroma/.
"""

# --- Importaciones estándar: rutas, fechas, hash, carga de archivos y logs. ---
import hashlib
import json
import os
import re
import shutil
import time
import urllib.error
import urllib.request
from datetime import datetime, timezone
from pathlib import Path

# --- Dependencias de la aplicación: interfaz, lectura de tablas y vector DB. ---
import chromadb
import pandas as pd
import streamlit as st
from chromadb.utils.embedding_functions import DefaultEmbeddingFunction

# --- Rutas del proyecto y extensiones que el pipeline puede procesar. ---
BASE_DIR = Path(__file__).parent
DOCS_DIR = BASE_DIR / "data" / "documents"
CHROMA_DIR = BASE_DIR / "data" / "chroma"
LOG_DIR = BASE_DIR / "logs"
SUPPORTED = {".pdf", ".docx", ".xlsx", ".xls", ".pptx", ".md", ".txt", ".csv", ".json", ".html", ".htm"}
COLLECTION_PREFIX = "documentos_corporativos"


def clean(text: str) -> str:
    """Normaliza espacios para evitar embeddings de texto vacío o con ruido."""
    return re.sub(r"\s+", " ", text or "").strip()


def markdown_to_text(text: str) -> str:
    """Elimina marcas de Markdown para que una respuesta no cree títulos enormes."""
    text = re.sub(r"!\[[^\]]*\]\([^)]*\)", "", text)  # imágenes
    text = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", text)  # enlaces
    text = re.sub(r"[*_`#>|]", "", text)  # negrita, encabezados y tablas
    return clean(text)


@st.cache_data(ttl=5, show_spinner=False)
def ollama_status() -> tuple[bool, str]:
    """Comprueba si el servidor local de Ollama está activo y expone modelos."""
    try:
        with urllib.request.urlopen("http://localhost:11434/api/tags", timeout=2) as response:
            models = json.loads(response.read().decode("utf-8")).get("models", [])
        return True, f"Ollama conectado · {len(models)} modelo(s) disponible(s)"
    except (urllib.error.URLError, urllib.error.HTTPError, TimeoutError, json.JSONDecodeError):
        return False, "Ollama no está disponible · se usará modo extractivo"


def metrics() -> dict:
    """Resume el log local sin usar servicios externos ni base de datos adicional."""
    log_file = LOG_DIR / "execution.jsonl"
    if not log_file.exists():
        return {"queries": 0, "no_answer": 0, "latency": 0}
    events = []
    for line in log_file.read_text(encoding="utf-8").splitlines():
        try:
            events.append(json.loads(line))
        except json.JSONDecodeError:
            continue
    if not events:
        return {"queries": 0, "no_answer": 0, "latency": 0}
    no_answer = sum("No encontré" in event.get("response", "") for event in events)
    return {"queries": len(events), "no_answer": no_answer, "latency": round(sum(event.get("latency_ms", 0) for event in events) / len(events))}


def category_for(path: Path) -> str:
    """Usa la primera carpeta bajo documents como metadato de categoría."""
    relative = path.relative_to(DOCS_DIR)
    return relative.parts[0] if len(relative.parts) > 1 else "Sin categoría"


def extract(path: Path) -> list[tuple[str, str]]:
    """Extrae pares (ubicación, texto) según el formato del archivo.

    Un error queda registrado como texto diagnóstico: un archivo defectuoso no
    impide que los demás documentos continúen siendo indexados.
    """
    ext = path.suffix.lower()
    try:
        if ext in {".md", ".txt"}:
            raw_text = path.read_text(encoding="utf-8", errors="ignore")
            return [("contenido", markdown_to_text(raw_text) if ext == ".md" else raw_text)]
        if ext == ".csv":
            return [("tabla", pd.read_csv(path).to_markdown(index=False))]
        if ext == ".json":
            data = json.loads(path.read_text(encoding="utf-8"))
            return [("estructura JSON", json.dumps(data, ensure_ascii=False, indent=2))]
        if ext in {".html", ".htm"}:
            from bs4 import BeautifulSoup
            html = path.read_text(encoding="utf-8", errors="ignore")
            return [("contenido HTML", BeautifulSoup(html, "html.parser").get_text(" "))]
        if ext == ".pdf":
            from pypdf import PdfReader
            return [(f"página {n}", page.extract_text() or "") for n, page in enumerate(PdfReader(path).pages, 1)]
        if ext == ".docx":
            from docx import Document
            return [("documento", "\n".join(p.text for p in Document(path).paragraphs))]
        if ext in {".xlsx", ".xls"}:
            book = pd.ExcelFile(path)
            return [(f"hoja {sheet}", pd.read_excel(path, sheet_name=sheet).to_markdown(index=False)) for sheet in book.sheet_names]
        if ext == ".pptx":
            from pptx import Presentation
            pres = Presentation(path)
            return [(f"diapositiva {n}", " ".join(s.text for s in slide.shapes if hasattr(s, "text"))) for n, slide in enumerate(pres.slides, 1)]
    except Exception as exc:
        return [("error de extracción", f"No se pudo leer el archivo: {exc}")]
    return []


def chunk(text: str, size: int = 750, overlap: int = 120) -> list[str]:
    """Divide el contenido en fragmentos solapados y prefiere cerrar en un punto."""
    text = clean(text)
    if not text:
        return []
    parts, start = [], 0
    while start < len(text):
        end = min(len(text), start + size)
        # Evita cortar una idea a la mitad cuando hay un fin de oración cercano.
        if end < len(text):
            boundary = text.rfind(". ", start, end)
            if boundary > start + size // 2:
                end = boundary + 1
        parts.append(text[start:end])
        # Al llegar al final no se debe aplicar overlap: hacerlo produciría
        # microfragmentos repetidos con las últimas palabras del documento.
        if end >= len(text):
            break
        start = max(end - overlap, start + 1)
    return parts


@st.cache_resource(show_spinner=False)
def vector_store(stamp: int):
    """Reconstruye el índice persistente cuando cambian los documentos.

    DefaultEmbeddingFunction usa all-MiniLM-L6-v2: transforma cada fragmento en
    un vector que representa su significado. Chroma guarda documentos, vectores
    y metadatos para reutilizarlos al reiniciar la aplicación.
    """
    CHROMA_DIR.mkdir(parents=True, exist_ok=True)
    client = chromadb.PersistentClient(path=str(CHROMA_DIR))
    # Cada reconstrucción recibe un nombre único. Así una ejecución de Streamlit
    # nunca elimina una colección que otra ejecución todavía está consultando.
    collection_name = f"{COLLECTION_PREFIX}_{stamp}"
    collection = client.create_collection(
        name=collection_name,
        embedding_function=DefaultEmbeddingFunction(),
        metadata={"hnsw:space": "cosine"},
    )

    ids, documents, metadatas = [], [], []
    for path in DOCS_DIR.rglob("*"):
        if not (path.is_file() and path.suffix.lower() in SUPPORTED):
            continue
        updated = datetime.fromtimestamp(path.stat().st_mtime, tz=timezone.utc).strftime("%Y-%m-%d")
        for location, text in extract(path):
            for number, content in enumerate(chunk(text)):
                # Un identificador estable evita duplicados al reindexar.
                raw_id = f"{path.relative_to(DOCS_DIR)}|{location}|{number}|{content}"
                ids.append(hashlib.sha256(raw_id.encode("utf-8")).hexdigest())
                documents.append(content)
                metadatas.append({"file": path.name, "category": category_for(path), "location": location, "updated": updated})
    # Inserción en lote: más eficiente que generar un embedding por llamada.
    if documents:
        collection.add(ids=ids, documents=documents, metadatas=metadatas)
    # Conserva los tres índices más recientes para no acumular archivos locales.
    collections = sorted(
        (getattr(item, "name", item) for item in client.list_collections()),
        reverse=True,
    )
    collections = [name for name in collections if name.startswith(COLLECTION_PREFIX)]
    for old_name in collections[3:]:
        try:
            client.delete_collection(old_name)
        except Exception:
            pass
    return collection, len(documents)


def retrieve(question: str, category: str, k: int = 4) -> list[dict]:
    """Busca semánticamente y reclasifica con coincidencias de términos clave."""
    collection, count = vector_store(st.session_state.index_stamp)
    if not count:
        return []
    # El filtro se aplica dentro de Chroma y no después de traer todos los datos.
    where = None if category == "Todas" else {"category": category}
    # Traemos más candidatos para que el reranker pueda corregir resultados
    # semánticamente parecidos pero poco específicos (por ejemplo, días y gastos).
    result = collection.query(query_texts=[question], n_results=min(max(k * 4, 12), count), where=where, include=["documents", "metadatas", "distances"])
    records = []
    # Palabras informativas de la pregunta. Se ignoran conectores frecuentes.
    ignored = {"que", "qué", "como", "cómo", "para", "por", "con", "una", "uno", "los", "las", "del", "tengo", "tiene", "cuanto", "cuántos", "cuantas", "cuántas", "dias", "días", "donde", "dónde", "puedo", "hacer"}
    keywords = {word for word in re.findall(r"\b\w+\b", question.lower()) if len(word) >= 4 and word not in ignored}
    for document, metadata, distance in zip(result["documents"][0], result["metadatas"][0], result["distances"][0]):
        searchable = f"{metadata['file']} {document}".lower()
        lexical_score = sum(word in searchable for word in keywords) / max(len(keywords), 1)
        semantic_score = max(0.0, 1 - float(distance))
        # El embedding aporta significado; la coincidencia exacta premia la política
        # que sí menciona el concepto solicitado.
        combined = 0.55 * semantic_score + 0.45 * lexical_score
        records.append({"text": document, **metadata, "score": combined, "semantic_score": semantic_score})
    return sorted(records, key=lambda record: record["score"], reverse=True)[:k]


def retrieval_question(question: str, messages: list[dict]) -> str:
    """Añade el último tema cuando una pregunta corta claramente es seguimiento."""
    follow_up = {"eso", "esa", "ese", "anterior", "también", "tambien", "y cuánto", "y cuanto"}
    previous_questions = [message["content"] for message in messages if message["role"] == "user"]
    if previous_questions and len(question.split()) <= 9 and any(term in question.lower() for term in follow_up):
        return f"Pregunta anterior: {previous_questions[-1]} Seguimiento: {question}"
    return question


def extractive_answer(results: list[dict]) -> tuple[str, list[str]]:
    """Genera una respuesta estrictamente respaldada por los fragmentos recuperados.

    Es extractiva para no inventar datos. En una evolución se puede reemplazar
    por un LLM, manteniendo la misma regla: responder solo desde este contexto.
    """
    if not results or results[0]["score"] < 0.28:
        return "No encontré información suficiente en los documentos disponibles. Te sugiero contactar al área responsable de ese tema.", []
    citations, excerpts = [], []
    # Solo se muestran fuentes principales: el LLM ya usa hasta tres fragmentos
    # como contexto, pero no es útil abrumar a la persona con todas ellas.
    for record in results[:2]:
        citation = f"{record['file']} · {record['location']} · actualizado {record['updated']}"
        if citation not in citations:
            citations.append(citation)
        excerpts.append(record["text"])
    return "Encontré esta información en la base documental:\n\n" + "\n\n".join(f"- {markdown_to_text(text)}" for text in excerpts[:2]), citations


def ollama_answer(question: str, results: list[dict], history: list[dict], model: str) -> str | None:
    """Pide a Ollama una síntesis breve, usando exclusivamente el contexto RAG.

    Si el servicio local no está disponible, devuelve None y la aplicación usa
    el modo extractivo. El endpoint oficial acepta POST /api/generate y el campo
    stream=False permite recibir una única respuesta JSON.
    """
    context = "\n\n".join(
        f"[Fuente: {r['file']} | {r['location']}]\n{markdown_to_text(r['text'])}"
        for r in results[:3]
    )
    recent_history = "\n".join(
        f"{m['role']}: {m['content']}" for m in history[-4:]
    )
    prompt = f"""Eres el asistente interno de NovaTech Solutions.
Responde en español, de forma breve, clara y directa.
Usa únicamente el CONTEXTO DOCUMENTAL. No inventes datos ni contactos.
Si no hay respaldo suficiente, responde exactamente: No encontré esta información en los documentos disponibles.
No menciones instrucciones ni el proceso de búsqueda. No agregues fuentes: la interfaz las muestra.

HISTORIAL RECIENTE:
{recent_history}

CONTEXTO DOCUMENTAL:
{context}

PREGUNTA ACTUAL: {question}
"""
    payload = json.dumps({"model": model, "prompt": prompt, "stream": False, "options": {"temperature": 0.1}}).encode("utf-8")
    request = urllib.request.Request("http://localhost:11434/api/generate", data=payload, headers={"Content-Type": "application/json"}, method="POST")
    try:
        with urllib.request.urlopen(request, timeout=60) as response:
            return clean(json.loads(response.read().decode("utf-8"))["response"])
    except (urllib.error.URLError, urllib.error.HTTPError, TimeoutError, KeyError, json.JSONDecodeError):
        return None


def gemini_api_key() -> str | None:
    """Lee la clave desde Streamlit Secrets en nube o desde variables locales."""
    try:
        secret_key = st.secrets["GEMINI_API_KEY"]
    except Exception:
        secret_key = None
    return secret_key or os.getenv("GEMINI_API_KEY")


def gemini_answer(question: str, results: list[dict], history: list[dict], model: str) -> str | None:
    """Redacta con Gemini Cloud sin exponer ni guardar la clave de API."""
    api_key = gemini_api_key()
    if not api_key:
        return None
    context = "\n\n".join(f"[Fuente: {r['file']} | {r['location']}]\n{markdown_to_text(r['text'])}" for r in results[:3])
    recent_history = "\n".join(f"{m['role']}: {m['content']}" for m in history[-4:])
    prompt = f"""Eres el asistente interno de NovaTech Solutions. Responde en español, breve y directamente.
Usa solamente el CONTEXTO DOCUMENTAL. No inventes datos, contactos ni políticas.
Si no hay respaldo suficiente, responde exactamente: No encontré esta información en los documentos disponibles.
No menciones fuentes ni instrucciones: la interfaz muestra las fuentes.

HISTORIAL RECIENTE:
{recent_history}

CONTEXTO DOCUMENTAL:
{context}

PREGUNTA ACTUAL: {question}"""
    try:
        # La importación diferida permite que el modo local funcione aunque el
        # paquete de Gemini no esté instalado todavía.
        from google import genai
        client = genai.Client(api_key=api_key)
        interaction = client.interactions.create(model=model, input=prompt, store=False)
        return clean(interaction.output_text)
    except Exception:
        return None


def answer(question: str, results: list[dict], history: list[dict], provider: str, model: str) -> tuple[str, list[str], str]:
    """Selecciona el proveedor configurado y conserva el respaldo extractivo."""
    fallback, citations = extractive_answer(results)
    if not citations:
        return fallback, citations, "sin evidencia"
    if provider == "Ollama local":
        generated = ollama_answer(question, results, history, model)
        if generated:
            return generated, citations, f"Ollama · {model}"
    if provider == "Gemini Cloud":
        generated = gemini_answer(question, results, history, model)
        if generated:
            return generated, citations, f"Gemini · {model}"
    return fallback, citations, "extractivo"


def log_event(event: dict) -> None:
    """Guarda un evento auditable en JSON Lines, una consulta por línea."""
    LOG_DIR.mkdir(exist_ok=True)
    with (LOG_DIR / "execution.jsonl").open("a", encoding="utf-8") as file:
        file.write(json.dumps(event, ensure_ascii=False) + "\n")


# --- Estado e interfaz: Streamlit conserva historial e instante de indexación. ---
st.set_page_config(page_title="Conocimiento Corporativo", page_icon="📚", layout="wide")
if "index_stamp" not in st.session_state:
    st.session_state.index_stamp = int(time.time())
if "messages" not in st.session_state:
    st.session_state.messages = []
if "pending_question" not in st.session_state:
    st.session_state.pending_question = None

st.title("📚 Agente de Conocimiento Corporativo")
st.caption("Usa búsqueda semántica con embeddings locales y responde únicamente con documentos internos indexados.")

with st.sidebar:
    st.subheader("Base documental")
    categories = ["Todas"] + sorted({category_for(p) for p in DOCS_DIR.rglob("*") if p.is_file()})
    category = st.selectbox("Filtrar por categoría", categories)
    st.divider()
    st.subheader("Redacción de respuestas")
    provider = st.radio("Proveedor", ["Extractivo (sin LLM)", "Ollama local", "Gemini Cloud"], help="Gemini Cloud funciona en un despliegue web; Ollama local funciona en tu computadora.")
    default_model = "llama3.2:3b" if provider == "Ollama local" else "gemini-3.5-flash"
    model_name = st.text_input("Modelo", value=default_model, disabled=provider == "Extractivo (sin LLM)", key=f"model_{provider}")
    if provider == "Ollama local":
        connected, status_message = ollama_status()
        (st.success if connected else st.info)(status_message)
    elif provider == "Gemini Cloud":
        if gemini_api_key():
            st.success("Gemini Cloud configurado")
        else:
            st.warning("Falta GEMINI_API_KEY en los secretos de Streamlit")
    if st.button("🧹 Nueva conversación", use_container_width=True):
        st.session_state.messages = []
        st.rerun()
    upload = st.file_uploader("Agregar documento", type=sorted(s[1:] for s in SUPPORTED))
    upload_category = st.selectbox("Categoría del archivo", [c for c in categories if c != "Todas"] or ["General"])
    if upload and st.button("Indexar archivo"):
        target = DOCS_DIR / upload_category / upload.name
        target.parent.mkdir(parents=True, exist_ok=True)
        with target.open("wb") as file:
            shutil.copyfileobj(upload, file)
        # time_ns evita reutilizar el mismo identificador si hay dos clics rápidos.
        st.session_state.index_stamp = time.time_ns()
        st.success(f"{upload.name} fue indexado con embeddings.")
    if st.button("Reconstruir índice"):
        st.session_state.index_stamp = time.time_ns()
        st.success("Índice vectorial actualizado.")
    _, count = vector_store(st.session_state.index_stamp)
    st.caption(f"{count} fragmentos vectorizados disponibles")
    with st.expander("Documentos indexados"):
        indexed_files = [{"Archivo": path.name, "Categoría": category_for(path), "Actualizado": datetime.fromtimestamp(path.stat().st_mtime).strftime("%Y-%m-%d")} for path in DOCS_DIR.rglob("*") if path.is_file() and path.suffix.lower() in SUPPORTED]
        st.dataframe(pd.DataFrame(indexed_files), hide_index=True, use_container_width=True)
    with st.expander("Métricas locales"):
        summary = metrics()
        first, second, third = st.columns(3)
        first.metric("Consultas", summary["queries"])
        second.metric("Sin respuesta", summary["no_answer"])
        third.metric("Latencia media", f"{summary['latency']} ms")

# Se vuelve a dibujar el historial completo en cada actualización de Streamlit.
for item in st.session_state.messages:
    with st.chat_message(item["role"]):
        st.markdown(item["content"])
        if item.get("citations"):
            st.caption("Fuentes: " + " | ".join(item["citations"]))

# Botones útiles para una demostración sin tener que inventar preguntas.
st.caption("Prueba rápida:")
suggestion_columns = st.columns(3)
suggestions = ["¿Cuántos días de vacaciones tengo?", "¿Qué hago ante una incidencia crítica?", "¿Qué hago si recibo phishing?"]
for column, suggestion in zip(suggestion_columns, suggestions):
    if column.button(suggestion, use_container_width=True):
        st.session_state.pending_question = suggestion
        st.rerun()

# El widget se crea siempre: si se usa una sugerencia, el campo de chat no
# desaparece y la persona puede continuar escribiendo la siguiente pregunta.
typed_question = st.chat_input("Ej.: ¿Cuántos días de vacaciones tengo?")
question = st.session_state.pending_question or typed_question
st.session_state.pending_question = None
if question:
    st.session_state.messages.append({"role": "user", "content": question})
    with st.chat_message("user"):
        st.markdown(question)
    started = time.perf_counter()
    results = retrieve(retrieval_question(question, st.session_state.messages), category)
    response, citations, response_mode = answer(question, results, st.session_state.messages, provider, model_name)
    latency_ms = round((time.perf_counter() - started) * 1000)
    with st.chat_message("assistant"):
        st.markdown(response)
        if citations:
            st.caption("Fuentes: " + " | ".join(citations))
        st.feedback("thumbs", key=f"feedback_{len(st.session_state.messages)}")
    st.session_state.messages.append({"role": "assistant", "content": response, "citations": citations})
    # El log permite auditar fuentes, calidad de búsqueda y tiempo de respuesta.
    log_event({"timestamp": datetime.now(timezone.utc).isoformat(), "question": question, "category_filter": category, "sources": citations, "scores": [round(r["score"], 4) for r in results], "response": response, "mode": response_mode, "latency_ms": latency_ms})
